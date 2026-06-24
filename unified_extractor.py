"""
unified_extractor.py
====================
Unified document extraction pipeline combining:
- Marker: For PDF files (best table extraction quality)
- Dedicated parsers: openpyxl(.xlsx), python-docx(.docx), python-pptx(.pptx), etc.
- Tesseract OCR: For image files and images embedded in documents
- Unstructured: Ultimate fallback for all other formats
"""

import os
import re
from knowledge_schema import KnowledgeEntry, KnowledgeBase, classify_knowledge_type


MARKER_EXTENSIONS = {'.pdf'}
SUPPORTED_EXTENSIONS = {
    '.doc', '.docx', '.odt', '.rtf', '.txt', '.md', '.rst', '.org', 
    '.html', '.htm', '.xml', '.json', '.ndjson',
    '.csv', '.xls', '.xlsx', '.tsv', '.ppt', '.pptx',
    '.eml', '.msg', '.epub', '.yaml', '.yml',
    '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.heic'
}


def is_pdf_encrypted(file_path: str) -> bool:
    try:
        import fitz
        doc = fitz.open(file_path)
        is_encrypted = doc.needs_pass
        doc.close()
        return is_encrypted
    except ImportError:
        print("[WARN] PyMuPDF not installed, cannot detect encrypted PDFs")
        return False
    except Exception as e:
        print("[WARN] Failed to check PDF encryption: {e}")
        return False


def extract_knowledge_from_file(file_path: str, company: str = "") -> KnowledgeBase:
    _, ext = os.path.splitext(file_path.lower())
    
    if ext in MARKER_EXTENSIONS:
        if is_pdf_encrypted(file_path):
            print("[Extractor] Skipping encrypted PDF: {file_path}")
            kb = KnowledgeBase(source_file=os.path.basename(file_path))
            kb.metadata["encrypted"] = True
            return kb
        
        print("[Extractor] Using Marker for PDF: {file_path}")
        kb = _extract_with_marker(file_path)
    elif ext in SUPPORTED_EXTENSIONS:
        print("[Extractor] Processing {ext}: {file_path}")
        kb = _extract_with_best_parser(file_path, ext)
    else:
        print("[Extractor] Unknown format {ext}, trying Unstructured fallback: {file_path}")
        kb = _extract_with_unstructured_fallback(file_path)
    
    if company:
        for entry in kb.entries:
            entry.metadata["company"] = company
    
    return kb


def _extract_with_marker(file_path: str) -> KnowledgeBase:
    from marker_extractor import extract_knowledge_from_pdf
    return extract_knowledge_from_pdf(file_path)


def _extract_with_best_parser(file_path: str, ext: str) -> KnowledgeBase:
    IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.heic'}
    
    if ext in IMAGE_EXTENSIONS:
        kb = _extract_with_tesseract(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.xlsx':
        kb = _extract_with_openpyxl(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.xls':
        kb = _extract_with_pandas(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.docx':
        kb = _extract_with_python_docx(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext in ['.doc', '.odt', '.rtf']:
        kb = _extract_with_document_generic(file_path, ext)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.pptx':
        kb = _extract_with_python_pptx_enhanced(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.ppt':
        kb = _extract_with_python_pptx_enhanced(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext in ['.eml', '.msg']:
        kb = _extract_with_email_parser(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext in ['.csv', '.tsv']:
        kb = _extract_with_csv_parser(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.json':
        kb = _extract_with_json_parser(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.ndjson':
        kb = _extract_with_ndjson_parser(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext in ['.xml', '.yaml', '.yml']:
        kb = _extract_with_data_parser(file_path, ext)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext in ['.html', '.htm']:
        kb = _extract_with_html_parser(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext in ['.md', '.rst', '.org']:
        kb = _extract_with_markup_parser(file_path, ext)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    if ext == '.epub':
        kb = _extract_with_epub_parser(file_path)
        if kb.entries:
            return kb
        return _extract_with_unstructured_fallback(file_path)
    
    return _extract_with_unstructured_fallback(file_path)


def _extract_with_unstructured_fallback(file_path: str) -> KnowledgeBase:
    try:
        from unstructured.partition.auto import partition
        elements = partition(file_path)
    except ImportError:
        print("[WARN] unstructured not installed, falling back to basic text extraction")
        return _fallback_text_extraction(file_path)
    except Exception as e:
        print("[WARN] Partition failed: {e}. Trying basic extraction.")
        return _fallback_text_extraction(file_path)
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    for element in elements:
        text = getattr(element, 'text', '')
        if not text or len(text.strip()) < 10:
            continue
        
        clean_text = _clean_garbled_text(text)
        if not clean_text or len(clean_text.strip()) < 10:
            continue
        
        element_type = getattr(element, 'type', '').lower()
        
        if element_type in ['title', 'heading', 'header']:
            title = clean_text.strip()
            content = ""
            k_type = "general"
        elif element_type == 'list':
            title = "List Items"
            content = clean_text.strip()
            k_type = "general"
        elif element_type == 'table':
            title = "Table Content"
            content = clean_text.strip()
            k_type = "general"
        else:
            title = clean_text.strip()[:50] + "..." if len(clean_text) > 50 else clean_text.strip()
            content = clean_text.strip()
            k_type = classify_knowledge_type(title, content)
        
        keywords = _extract_keywords(clean_text)
        
        entry = KnowledgeEntry(
            type=k_type,
            title=title,
            content=content,
            source_file=os.path.basename(file_path),
            source_page=0,
            keywords=keywords,
            metadata={"element_type": element_type, "parser": "unstructured"}
        )
        kb.add(entry)
    
    print("[Unstructured] Extracted {len(kb.entries)} entries")
    return kb


def _extract_with_tesseract(file_path: str) -> KnowledgeBase:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        print("[WARN] pytesseract or Pillow not installed. Falling back to basic extraction.")
        return _fallback_text_extraction(file_path)
    
    import platform
    if platform.system() == "Windows":
        possible_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Users\{}\AppData\Local\Programs\Tesseract-OCR\tesseract.exe".format(os.getenv("USERNAME", ""))
        ]
        for p in possible_paths:
            if os.path.exists(p):
                pytesseract.pytesseract.tesseract_cmd = p
                break
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        img = Image.open(file_path)
        
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        custom_config = r'--oem 3 --psm 3 -l eng'
        full_text = pytesseract.image_to_string(img, config=custom_config)
        
        if not full_text or len(full_text.strip()) < 10:
            print("[Tesseract] No meaningful text extracted from {file_path}")
            return kb
        
        full_text = _clean_garbled_text(full_text)
        if not full_text or len(full_text.strip()) < 10:
            return kb
        
        paragraphs = [p.strip() for p in re.split(r'\n{2,}', full_text) if p.strip() and len(p.strip()) >= 15]
        
        if not paragraphs:
            paragraphs = [p.strip() for p in full_text.split('\n') if p.strip() and len(p.strip()) >= 15]
        
        for para in paragraphs:
            alpha_ratio = sum(1 for c in para if c.isalpha()) / max(len(para), 1)
            if alpha_ratio < 0.3:
                continue
            
            title = para[:50] + "..." if len(para) > 50 else para
            k_type = classify_knowledge_type(title, para)
            keywords = _extract_keywords(para)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=para,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "image_ocr", "ocr_engine": "tesseract"}
            )
            kb.add(entry)
        
        print("[Tesseract] Extracted {len(kb.entries)} entries from image")
    
    except pytesseract.TesseractNotFoundError:
        print("[ERROR] Tesseract executable not found.")
        return _fallback_text_extraction(file_path)
    except Exception as e:
        print("[ERROR] Tesseract OCR failed: {e}")
        return _fallback_text_extraction(file_path)
    
    return kb


def _clean_garbled_text(text: str) -> str:
    if not text:
        return ""
    
    printable = sum(1 for c in text if c.isprintable() or c in '\n\t\r ')
    if len(text) == 0:
        return ""
    
    printable_ratio = printable / len(text)
    
    if printable_ratio < 0.6:
        return ""
    
    text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f]', '', text)
    
    return text


def _fallback_text_extraction(file_path: str) -> KnowledgeBase:
    _, ext = os.path.splitext(file_path.lower())
    
    try:
        if ext in ['.txt', '.md', '.json', '.xml', '.rst', '.org', '.yaml', '.yml', '.ndjson']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif ext in ['.csv', '.tsv']:
            import csv
            text = ""
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " | ".join(row) + "\n"
        elif ext in ['.html', '.htm']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()
        else:
            with open(file_path, 'rb') as f:
                raw = f.read()
            text = raw.decode('utf-8', errors='ignore')[:50000]
    except Exception as e:
        print("[ERROR] Failed to read file: {e}")
        text = ""
    
    if not text:
        kb = KnowledgeBase(source_file=os.path.basename(file_path))
        return kb
    
    paragraphs = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip() and len(p.strip()) > 10]
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    for i, para in enumerate(paragraphs):
        title = para[:50] + "..." if len(para) > 50 else para
        k_type = classify_knowledge_type(title, para)
        keywords = _extract_keywords(para)
        
        entry = KnowledgeEntry(
            type=k_type,
            title=title,
            content=para,
            source_file=os.path.basename(file_path),
            source_page=0,
            keywords=keywords,
            metadata={"element_type": "text", "parser": "fallback"}
        )
        kb.add(entry)
    
    return kb


def _extract_with_openpyxl(file_path: str) -> KnowledgeBase:
    try:
        import openpyxl
    except ImportError:
        print("[WARN] openpyxl not installed")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            max_row = ws.max_row
            max_col = ws.max_column
            
            if max_row == 0 or max_col == 0:
                continue
            
            header_row = []
            for col in range(1, max_col + 1):
                cell = ws.cell(row=1, column=col)
                header_row.append(str(cell.value).strip() if cell.value else "")
            
            table_data = []
            for row in range(2, max_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell = ws.cell(row=row, column=col)
                    row_data.append(str(cell.value).strip() if cell.value else "")
                if any(row_data):
                    table_data.append(row_data)
            
            if not table_data:
                continue
            
            content_lines = ["Sheet: {sheet_name}"]
            if any(header_row):
                content_lines.append("Columns: {', '.join(header_row)}")
            for row in table_data[:100]:
                content_lines.append(" | ".join(row))
            
            content = "\n".join(content_lines)
            title = "Excel Table - {sheet_name}"
            k_type = "general"
            keywords = _extract_keywords(content)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=content,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "excel_table", "sheet": sheet_name, "parser": "openpyxl"}
            )
            kb.add(entry)
        
        print("[openpyxl] Extracted {len(kb.entries)} entries from Excel")
    
    except Exception as e:
        print("[ERROR] openpyxl extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_pandas(file_path: str) -> KnowledgeBase:
    try:
        import pandas as pd
    except ImportError:
        print("[WARN] pandas not installed")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        df = pd.read_excel(file_path, engine='xlrd')
        
        title = "Excel Data - {os.path.basename(file_path)}"
        content = df.to_string(index=False)[:10000]
        k_type = "general"
        keywords = _extract_keywords(content)
        
        entry = KnowledgeEntry(
            type=k_type,
            title=title,
            content=content,
            source_file=os.path.basename(file_path),
            source_page=0,
            keywords=keywords,
            metadata={"element_type": "excel_table", "parser": "pandas"}
        )
        kb.add(entry)
        
        print("[pandas] Extracted {len(kb.entries)} entries from Excel")
    
    except Exception as e:
        print("[ERROR] pandas extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_python_docx(file_path: str) -> KnowledgeBase:
    try:
        from docx import Document
    except ImportError:
        print("[WARN] python-docx not installed")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        doc = Document(file_path)
        
        for section in doc.sections:
            if section.header:
                for para in section.header.paragraphs:
                    text = para.text.strip()
                    if text and len(text) >= 10:
                        clean_text = _clean_garbled_text(text)
                        if clean_text:
                            title = "Header: " + clean_text[:30]
                            k_type = "general"
                            keywords = _extract_keywords(clean_text)
                            entry = KnowledgeEntry(
                                type=k_type,
                                title=title,
                                content=clean_text,
                                source_file=os.path.basename(file_path),
                                source_page=0,
                                keywords=keywords,
                                metadata={"element_type": "header", "parser": "python-docx"}
                            )
                            kb.add(entry)
            
            if section.footer:
                for para in section.footer.paragraphs:
                    text = para.text.strip()
                    if text and len(text) >= 10:
                        clean_text = _clean_garbled_text(text)
                        if clean_text:
                            title = "Footer: " + clean_text[:30]
                            k_type = "general"
                            keywords = _extract_keywords(clean_text)
                            entry = KnowledgeEntry(
                                type=k_type,
                                title=title,
                                content=clean_text,
                                source_file=os.path.basename(file_path),
                                source_page=0,
                                keywords=keywords,
                                metadata={"element_type": "footer", "parser": "python-docx"}
                            )
                            kb.add(entry)
        
        for i, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if not text or len(text) < 10:
                continue
            
            clean_text = _clean_garbled_text(text)
            if not clean_text:
                continue
            
            style_name = paragraph.style.name if paragraph.style else ""
            if 'Heading' in style_name or 'Title' in style_name:
                k_type = "general"
            else:
                k_type = classify_knowledge_type(text[:50], text)
            
            title = text[:50] + "..." if len(text) > 50 else text
            keywords = _extract_keywords(text)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=clean_text,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "paragraph", "style": style_name, "parser": "python-docx"}
            )
            kb.add(entry)
        
        for table in doc.tables:
            table_content = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_content.append(" | ".join(cells))
            
            content = "\n".join(table_content)
            if len(content) < 10:
                continue
            
            title = "Table Content"
            k_type = "general"
            keywords = _extract_keywords(content)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=content,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "table", "parser": "python-docx"}
            )
            kb.add(entry)
        
        print("[python-docx] Extracted {len(kb.entries)} entries from Word")
    
    except Exception as e:
        print("[ERROR] python-docx extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_document_generic(file_path: str, ext: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    if ext == '.odt':
        try:
            from odf.opendocument import load
            from odf.text import P, H, List
            doc = load(file_path)
            
            for element in doc.getElementsByType(P):
                text = str(element).strip()
                if text and len(text) >= 10:
                    clean_text = _clean_garbled_text(text)
                    if clean_text:
                        title = clean_text[:50] + "..." if len(clean_text) > 50 else clean_text
                        k_type = classify_knowledge_type(title, clean_text)
                        keywords = _extract_keywords(clean_text)
                        entry = KnowledgeEntry(
                            type=k_type,
                            title=title,
                            content=clean_text,
                            source_file=os.path.basename(file_path),
                            source_page=0,
                            keywords=keywords,
                            metadata={"element_type": "paragraph", "parser": "odfpy"}
                        )
                        kb.add(entry)
            
            print("[odfpy] Extracted {len(kb.entries)} entries from ODT")
            return kb
        except ImportError:
            print("[WARN] odfpy not installed")
            return _fallback_text_extraction(file_path)
        except Exception as e:
            print("[ERROR] ODT extraction failed: {e}")
            return _fallback_text_extraction(file_path)
    
    return _fallback_text_extraction(file_path)


def _extract_with_python_pptx_enhanced(file_path: str) -> KnowledgeBase:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("[WARN] python-pptx not installed")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        prs = Presentation(file_path)
        
        def extract_text_from_shape(shape):
            text = ""
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    text += extract_text_from_shape(subshape) + "\n"
            elif hasattr(shape, 'text'):
                text += shape.text + "\n"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    import pytesseract
                    from PIL import Image
                    import io
                    image_data = shape.image.blob
                    img = Image.open(io.BytesIO(image_data))
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    custom_config = r'--oem 3 --psm 3 -l eng'
                    ocr_text = pytesseract.image_to_string(img, config=custom_config)
                    if ocr_text and len(ocr_text.strip()) >= 5:
                        text += "[OCR] " + ocr_text + "\n"
                except Exception:
                    pass
            return text.strip()
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            
            for shape in slide.shapes:
                text = extract_text_from_shape(shape)
                if text and len(text) > 5:
                    slide_content.append(text)
            
            if not slide_content:
                continue
            
            content = "\n".join(slide_content)
            title = "Slide {slide_idx + 1}"
            k_type = "general"
            keywords = _extract_keywords(content)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=content,
                source_file=os.path.basename(file_path),
                source_page=slide_idx + 1,
                keywords=keywords,
                metadata={"element_type": "slide", "parser": "python-pptx-enhanced"}
            )
            kb.add(entry)
        
        print("[python-pptx] Extracted {len(kb.entries)} entries from PowerPoint")
    
    except Exception as e:
        print("[ERROR] python-pptx extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_email_parser(file_path: str) -> KnowledgeBase:
    _, ext = os.path.splitext(file_path.lower())
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        if ext == '.eml':
            import email
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                msg = email.message_from_file(f)
            
            subject = msg['subject'] or "No Subject"
            from_addr = msg['from'] or "Unknown Sender"
            to_addr = msg['to'] or "Unknown Recipient"
            date = msg['date'] or "Unknown Date"
            
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        payload = part.get_payload(decode=True)
                        if payload:
                            body = payload.decode('utf-8', errors='ignore')
                        break
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    body = payload.decode('utf-8', errors='ignore')
            
            content = "From: {from_addr}\nTo: {to_addr}\nDate: {date}\nSubject: {subject}\n\n{body}"
            title = "Email: {subject}"
            k_type = "general"
            keywords = _extract_keywords(content)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=content,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "email", "parser": "email"}
            )
            kb.add(entry)
        
        elif ext == '.msg':
            try:
                import extract_msg
                msg = extract_msg.Message(file_path)
                
                content = "From: {msg.sender}\nTo: {msg.to}\nDate: {msg.date}\nSubject: {msg.subject}\n\n{msg.body}"
                title = "Email: {msg.subject}"
                k_type = "general"
                keywords = _extract_keywords(content)
                
                entry = KnowledgeEntry(
                    type=k_type,
                    title=title,
                    content=content,
                    source_file=os.path.basename(file_path),
                    source_page=0,
                    keywords=keywords,
                    metadata={"element_type": "email", "parser": "extract_msg"}
                )
                kb.add(entry)
            except ImportError:
                print("[WARN] extract-msg not installed")
                return _fallback_text_extraction(file_path)
        
        print("[email] Extracted {len(kb.entries)} entries from email")
    
    except Exception as e:
        print("[ERROR] email extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_csv_parser(file_path: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        import csv
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            rows = list(reader)
        
        if not rows:
            return kb
        
        header = rows[0]
        data_rows = rows[1:]
        
        content_lines = ["Columns: {', '.join(header)}"]
        for row in data_rows[:200]:
            content_lines.append(" | ".join(row))
        
        content = "\n".join(content_lines)
        title = "CSV Table - {os.path.basename(file_path)}"
        k_type = "general"
        keywords = _extract_keywords(content)
        
        entry = KnowledgeEntry(
            type=k_type,
            title=title,
            content=content,
            source_file=os.path.basename(file_path),
            source_page=0,
            keywords=keywords,
            metadata={"element_type": "csv_table", "parser": "csv"}
        )
        kb.add(entry)
        
        print("[csv] Extracted {len(kb.entries)} entries from CSV")
    
    except Exception as e:
        print("[ERROR] CSV extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_json_parser(file_path: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        import json
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
        
        content = json.dumps(data, indent=2, ensure_ascii=False)[:50000]
        title = "JSON Data - {os.path.basename(file_path)}"
        k_type = "general"
        keywords = _extract_keywords(content)
        
        entry = KnowledgeEntry(
            type=k_type,
            title=title,
            content=content,
            source_file=os.path.basename(file_path),
            source_page=0,
            keywords=keywords,
            metadata={"element_type": "json", "parser": "json"}
        )
        kb.add(entry)
        
        print("[json] Extracted {len(kb.entries)} entries from JSON")
    
    except Exception as e:
        print("[ERROR] JSON extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_ndjson_parser(file_path: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        import json
        
        entries_data = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            for line_num, line in enumerate(f):
                line = line.strip()
                if line:
                    try:
                        data = json.loads(line)
                        entries_data.append(data)
                    except json.JSONDecodeError:
                        pass
        
        if not entries_data:
            return kb
        
        content = json.dumps(entries_data, indent=2, ensure_ascii=False)[:50000]
        title = "NDJSON Data - {os.path.basename(file_path)}"
        k_type = "general"
        keywords = _extract_keywords(content)
        
        entry = KnowledgeEntry(
            type=k_type,
            title=title,
            content=content,
            source_file=os.path.basename(file_path),
            source_page=0,
            keywords=keywords,
            metadata={"element_type": "ndjson", "parser": "ndjson"}
        )
        kb.add(entry)
        
        print("[ndjson] Extracted {len(kb.entries)} entries from NDJSON")
    
    except Exception as e:
        print("[ERROR] NDJSON extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_data_parser(file_path: str, ext: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        if ext == '.xml':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()
        elif ext in ['.yaml', '.yml']:
            try:
                import yaml
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    data = yaml.safe_load(f)
                text = json.dumps(data, indent=2, ensure_ascii=False)[:50000]
            except ImportError:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
        
        if not text or len(text) < 10:
            return kb
        
        paragraphs = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip() and len(p.strip()) >= 15]
        
        for para in paragraphs:
            title = para[:50] + "..." if len(para) > 50 else para
            k_type = classify_knowledge_type(title, para)
            keywords = _extract_keywords(para)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=para,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": ext.strip('.'), "parser": "data"}
            )
            kb.add(entry)
        
        print("[data] Extracted {len(kb.entries)} entries from {ext}")
    
    except Exception as e:
        print("[ERROR] {ext} extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_html_parser(file_path: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f, 'html.parser')
            
            for script in soup(['script', 'style']):
                script.decompose()
            
            text = soup.get_text()
            text = re.sub(r'\s+', ' ', text).strip()
        except ImportError:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()
        
        if not text or len(text) < 10:
            return kb
        
        paragraphs = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip() and len(p.strip()) >= 15]
        
        for para in paragraphs:
            title = para[:50] + "..." if len(para) > 50 else para
            k_type = classify_knowledge_type(title, para)
            keywords = _extract_keywords(para)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=para,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "html", "parser": "html"}
            )
            kb.add(entry)
        
        print("[html] Extracted {len(kb.entries)} entries from HTML")
    
    except Exception as e:
        print("[ERROR] HTML extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_markup_parser(file_path: str, ext: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        if not text or len(text) < 10:
            return kb
        
        paragraphs = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip() and len(p.strip()) >= 10]
        
        for para in paragraphs:
            title = para[:50] + "..." if len(para) > 50 else para
            k_type = classify_knowledge_type(title, para)
            keywords = _extract_keywords(para)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=para,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": ext.strip('.'), "parser": "markup"}
            )
            kb.add(entry)
        
        print("[markup] Extracted {len(kb.entries)} entries from {ext}")
    
    except Exception as e:
        print("[ERROR] {ext} extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_with_epub_parser(file_path: str) -> KnowledgeBase:
    kb = KnowledgeBase(source_file=os.path.basename(file_path))
    
    try:
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(file_path)
            
            all_text = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.content, 'html.parser')
                    text = soup.get_text()
                    text = re.sub(r'\s+', ' ', text).strip()
                    if text and len(text) > 20:
                        all_text.append(text)
            
            if not all_text:
                return kb
            
            content = "\n\n".join(all_text)[:50000]
            title = "EPUB Book - {book.title if book.title else os.path.basename(file_path)}"
            k_type = "general"
            keywords = _extract_keywords(content)
            
            entry = KnowledgeEntry(
                type=k_type,
                title=title,
                content=content,
                source_file=os.path.basename(file_path),
                source_page=0,
                keywords=keywords,
                metadata={"element_type": "epub", "parser": "ebooklib"}
            )
            kb.add(entry)
            
            print("[epub] Extracted {len(kb.entries)} entries from EPUB")
        except ImportError:
            print("[WARN] ebooklib not installed")
            return _fallback_text_extraction(file_path)
    
    except Exception as e:
        print("[ERROR] EPUB extraction failed: {e}")
        return KnowledgeBase(source_file=os.path.basename(file_path))
    
    return kb


def _extract_keywords(text: str, max_kw: int = 8) -> list:
    keywords = []
    clean_text = re.sub(r'[<>{}[\]\\|]', ' ', text)
    
    upper_words = re.findall(r'\b[A-Z]{3,}\b', clean_text)
    keywords.extend(list(set(upper_words))[:4])
    
    title_words = re.findall(r'(?:^|\n|\.\s|\:\s)([A-Z][a-zA-Z]{2,})', clean_text)
    stop_words = {'the', 'and', 'for', 'are', 'you', 'your', 'this', 'that', 'with', 'from', 
                  'into', 'onto', 'over', 'under', 'through', 'between', 'among', 'without', 
                  'within', 'upon', 'than', 'but', 'or', 'so', 'yet', 'because', 'while', 'if', 
                  'when', 'though', 'although', 'unless', 'until', 'since', 'as', 'than', 'that', 
                  'which', 'who', 'whom', 'whose', 'what', 'where', 'why', 'how', 'is', 'are', 
                  'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 
                  'will', 'would', 'could', 'should', 'may', 'might', 'must', 'shall', 'can'}
    keywords.extend([w for w in title_words if w.lower() not in stop_words][:4])
    
    step_patterns = re.findall(r'\b(Step\s+\d+|Phase\s+\d+|Chapter\s+\d+|Section\s+\d+)\b', clean_text, re.IGNORECASE)
    keywords.extend(step_patterns[:2])
    
    tech_terms = re.findall(r'\b(printer|scanner|document|network|wireless|USB|Wi-Fi|Ethernet|software|driver|settings|paper|cartridge|ink|toner|control|panel|display|button|menu|option|feature|function|system|device|installation|setup|configuration|connection|interface|memory|storage|data|file|folder|image|text|table|page|format|size|quality|speed|performance|support|help|guide|manual|troubleshoot|error|problem|solution|warning|caution|note)\b', clean_text, re.IGNORECASE)
    keywords.extend([w.lower() for w in tech_terms][:3])
    
    seen = set()
    result = []
    for kw in keywords:
        kw_clean = kw.strip()
        if len(kw_clean) >= 3 and kw_clean.lower() not in stop_words:
            kw_normalized = kw_clean.lower()
            if kw_normalized not in seen:
                seen.add(kw_normalized)
                result.append(kw_clean)
        if len(result) >= max_kw:
            break
    
    return result